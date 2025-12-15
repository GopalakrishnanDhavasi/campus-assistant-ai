import fitz  # PyMuPDF
import pdfplumber
import re
import time
import traceback
from typing import List, Dict, Any, Optional
from sentence_transformers import SentenceTransformer
from .config import llm_client
import docx
from pptx import Presentation
import os
import json
import glob
import wikipedia

# --- LLM HELPER FUNCTIONS ---

def _extract_model_names(models_resp: Any) -> List[str]:
    names: List[str] = []
    data = getattr(models_resp, "data", None) or (models_resp if isinstance(models_resp, list) else None)
    if not data:
        return names
    for item in data:
        if isinstance(item, dict):
            mid = item.get("id") or item.get("model") or item.get("name")
        else:
            mid = getattr(item, "id", None) or getattr(item, "model", None) or getattr(item, "name", None)
        if mid:
            names.append(str(mid))
    return names

def pick_fallback_model(preferred: Optional[str] = None) -> Optional[str]:
    try:
        models_resp = llm_client.models.list()
        names = _extract_model_names(models_resp)
        if not names:
            return preferred or None
        if preferred and preferred in names:
            return preferred
        families = ("llama-3.1", "llama-3", "llama3", "llama", "gpt-oss", "gpt", "mixtral", "gemma")
        for fam in families:
            for n in names:
                if fam in n:
                    return n
        return names[0]
    except Exception:
        return preferred or None

def _extract_text_from_response(resp: Any) -> Optional[str]:
    try:
        out = getattr(resp, "output_text", None)
        if out: return str(out).strip()
        if isinstance(resp, dict):
            if "output_text" in resp and resp["output_text"]:
                return str(resp["output_text"]).strip()
            out_list = resp.get("output") or resp.get("choices") or None
            if out_list and isinstance(out_list, list):
                texts = []
                for item in out_list:
                    if isinstance(item, dict):
                        content = item.get("content")
                        if isinstance(content, list):
                            for c in content:
                                if isinstance(c, dict) and "text" in c: texts.append(str(c["text"]))
                        if "text" in item and isinstance(item["text"], str): texts.append(item["text"])
                        message = item.get("message")
                        if isinstance(message, dict):
                            msg_text = message.get("content")
                            if isinstance(msg_text, str): texts.append(msg_text)
                if texts: return "\n".join(texts).strip()
        out_attr = getattr(resp, "output", None)
        if out_attr and isinstance(out_attr, list):
            parts = []
            for entry in out_attr:
                if isinstance(entry, dict):
                    content = entry.get("content")
                    if isinstance(content, str): parts.append(content)
                else:
                    content = getattr(entry, "content", None)
                    if isinstance(content, list):
                        for c in content:
                            txt = c.get("text") if isinstance(c, dict) else getattr(c, "text", None)
                            if txt: parts.append(txt)
            if parts: return " ".join(parts).strip()
    except Exception:
        pass
    try: return str(resp).strip()
    except Exception: return None

def get_wikipedia_summary(query: str) -> str:
    """
    Searches Wikipedia for the query and returns a summary.
    Returns None if no page is found or an error occurs.
    """
    try:
        # 1. Search for the best matching page title
        search_results = wikipedia.search(query)
        if not search_results:
            return None
        
        # 2. Get the summary of the top result (first 3 sentences)
        summary = wikipedia.summary(search_results[0], sentences=3)
        return f"{summary}\n\n(Source: Wikipedia - {search_results[0]})"
    except wikipedia.exceptions.DisambiguationError as e:
        # If the term is too vague (e.g., "Python"), pick the first option
        try:
            summary = wikipedia.summary(e.options[0], sentences=3)
            return f"{summary}\n\n(Source: Wikipedia - {e.options[0]})"
        except:
            return None
    except Exception as e:
        print(f"Wikipedia Error: {e}")
        return None

def call_llm_summarize(prompt: str, max_tokens: int = 1024, model: Optional[str] = None, temperature: float = 0.0, stop: Optional[List[str]] = None, retry: int = 1, backoff_base: float = 1.0) -> Dict[str, Any]:
    chosen = pick_fallback_model(model)
    last_tb = None
    for attempt in range(0, retry + 1):
        try:
            call_kwargs: Dict[str, Any] = {"model": chosen, "input": prompt, "max_output_tokens": max_tokens, "temperature": temperature}
            if stop: call_kwargs["stop"] = stop
            resp = llm_client.responses.create(**call_kwargs)
            text = _extract_text_from_response(resp)
            return {"ok": True, "text": text, "resp": resp}
        except Exception as e:
            last_tb = traceback.format_exc()
            err_txt = str(e).lower()
            if any(token in err_txt for token in ("model", "not found", "decommissioned", "invalid_request", "unsupported")):
                chosen = pick_fallback_model(None)
                time.sleep(backoff_base * (2 ** attempt))
                continue
            if "rate_limit" in err_txt or "tokens" in err_txt or "too large" in err_txt:
                time.sleep(backoff_base * (2 ** attempt))
                continue
            return {"ok": False, "text": None, "resp": last_tb}
    return {"ok": False, "text": None, "resp": last_tb}

def call_llm_text_only(prompt: str, max_tokens: int, temperature: float) -> str:
    result = call_llm_summarize(prompt, max_tokens=max_tokens, temperature=temperature)
    if result.get("ok") and result.get("text"):
        return result["text"]
    return ""

def call_llm_answer(prompt_text: str, max_tokens: int = 512, temperature: float = 0.0):
    out = call_llm_summarize(prompt_text, max_tokens=max_tokens, temperature=temperature)
    if isinstance(out, dict) and "text" in out:
        return out["text"]
    return str(out)

# ... existing imports ...

def generate_multi_queries(original_question: str, call_llm_fn, n_versions: int = 3) -> List[str]:
    """
    Asks the LLM to generate 'n_versions' of the user's question
    from different perspectives to improve retrieval coverage.
    """
    prompt = f"""
    You are an AI language model assistant. Your task is to generate {n_versions} different versions of the given user question to retrieve relevant documents from a vector database. 
    By generating multiple perspectives on the user question, your goal is to help the user overcome some of the limitations of the distance-based similarity search. 
    
    Original question: {original_question}
    
    Provide these alternative questions separated by newlines. Do not number them.
    """
    
    # 1. Get response from LLM
    # We use a slightly higher temperature (0.7) to encourage creativity/variety
    response_text = call_llm_fn(prompt, max_tokens=256, temperature=0.7)
    
    # 2. Clean and split the response into a list
    variations = [line.strip() for line in response_text.split('\n') if line.strip()]
    
    # Return the original question + the new variations
    return [original_question] + variations

# --- PDF & DB HELPER FUNCTIONS ---

# --- IMPORTS FOR NEW FORMATS ---
import docx
from pptx import Presentation

# ... (Keep your existing Imports and LLM Helper Functions at the top) ...

# --- UNIVERSAL FILE PROCESSING FUNCTIONS ---

def extract_text_from_pdf_selectable(path: str) -> List[Dict[str, Any]]:
    # ... (Keep your existing PDF extraction code exactly as is) ...
    pages = []
    try:
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            text = page.get_text("text") or ""
            pages.append({"page_number": i + 1, "pdf_text": text})
        doc.close()
        # Fallback logic...
        total_chars = sum(len(p["pdf_text"]) for p in pages)
        if total_chars < 50:
            raise ValueError("PyMuPDF returned very little text")
        return pages
    except Exception:
        # Fallback to pdfplumber
        pages = []
        with pdfplumber.open(path) as pdf:
            for i, p in enumerate(pdf.pages):
                try: text = p.extract_text() or ""
                except Exception: text = ""
                pages.append({"page_number": i + 1, "pdf_text": text})
        return pages

def extract_text_from_docx(path: str) -> List[Dict[str, Any]]:
    """Reads Word files. Treats the whole document as 'Page 1' for simplicity."""
    try:
        doc = docx.Document(path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        text = "\n".join(full_text)
        return [{"page_number": 1, "pdf_text": text}]
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return []

def extract_text_from_pptx(path: str) -> List[Dict[str, Any]]:
    """Reads PPTX files. Maps Slides to Pages."""
    try:
        prs = Presentation(path)
        pages = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            text = "\n".join(text_runs)
            pages.append({"page_number": i + 1, "pdf_text": text})
        return pages
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return []

def extract_text_universal(path: str) -> List[Dict[str, Any]]:
    """Router: Checks file extension and calls the right extractor."""
    ext = os.path.splitext(path)[1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf_selectable(path)
    elif ext == ".docx":
        return extract_text_from_docx(path)
    elif ext == ".pptx":
        return extract_text_from_pptx(path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def preprocess_for_llm(text: str) -> str:
    text = re.sub(r'(\w)-\n(\w)', r'\1\2', text)
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
    text = re.sub(r'\s{2,}', ' ', text)
    return text.strip()

def build_combined_document(pages: List[Dict[str, Any]]) -> Dict[str, Any]:
    per_page = []
    for p in pages:
        cleaned = preprocess_for_llm(p["pdf_text"] or "")
        per_page.append({"page_number": p["page_number"], "cleaned_text": cleaned})
    combined_text = "\n\n".join([f"--- PAGE {p['page_number']} ---\n{p['cleaned_text']}" for p in per_page])
    return {"combined_text": combined_text, "per_page": per_page}

def simple_chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
    chunks = []
    start = 0
    text_len = len(text)
    chunk_id = 0
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunk_text = text[start:end].strip()
        if chunk_text:
            chunks.append({"id": f"chunk_{chunk_id}", "text": chunk_text, "start_char": start, "end_char": end})
            chunk_id += 1
        start = end - overlap
        if start < 0: start = 0
        if end >= text_len: break
    return chunks

def upsert_chunks_to_chroma(chunks: List[Dict[str, Any]], embed_model: SentenceTransformer, collection):
    texts = [c["text"] for c in chunks]
    ids = [c["id"] for c in chunks]
    metadatas = [{"start_char": c["start_char"], "end_char": c["end_char"]} for c in chunks]
    embeddings = embed_model.encode(texts, convert_to_numpy=True, show_progress_bar=True)
    emb_list = [emb.tolist() for emb in embeddings]
    collection.add(ids=ids, documents=texts, metadatas=metadatas, embeddings=emb_list)
    return collection

def get_all_chunks_from_collection(collection):
    try:
        res = collection.get(include=["documents", "metadatas", "embeddings"])
        ids = res.get("ids")
        docs = res.get("documents", [])
        metas = res.get("metadatas", [])
        if ids and isinstance(ids[0], list): ids = ids[0]
        if docs and isinstance(docs[0], list): docs = docs[0]
        if metas and isinstance(metas[0], list): metas = metas[0]
        chunks = []
        if not ids: ids = [f"doc_{i}" for i in range(len(docs))]
        for i in range(len(docs)):
            chunk_id = ids[i] if i < len(ids) else f"doc_{i}"
            metadata = metas[i] if i < len(metas) else {}
            chunks.append({"id": chunk_id, "text": docs[i], "metadata": metadata})
        return chunks
    except Exception as e:
        print("Warning: collection.get() failed:", str(e))
        return []

# --- SUMMARY STORAGE HELPERS ---
SUMMARY_STORAGE_DIR = "saved_summaries"
os.makedirs(SUMMARY_STORAGE_DIR, exist_ok=True)

def save_summary_to_disk(filename: str, summary_text: str):
    """Saves a summary to a text file."""
    safe_name = "".join([c for c in filename if c.isalpha() or c.isdigit() or c in " ._-"])
    file_path = os.path.join(SUMMARY_STORAGE_DIR, f"{safe_name}.txt")
    
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(summary_text)
    return safe_name

def list_saved_summaries():
    """Returns a list of saved summary filenames."""
    files = glob.glob(os.path.join(SUMMARY_STORAGE_DIR, "*.txt"))
    return [os.path.splitext(os.path.basename(f))[0] for f in files]

def load_summary_from_disk(filename: str):
    """Loads a summary by filename."""
    safe_name = "".join([c for c in filename if c.isalpha() or c.isdigit() or c in " ._-"])
    file_path = os.path.join(SUMMARY_STORAGE_DIR, f"{safe_name}.txt")
    
    if not os.path.exists(file_path):
        return None
        
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()
    
QUIZ_STORAGE_DIR = "saved_quizzes"
os.makedirs(QUIZ_STORAGE_DIR, exist_ok=True)

def save_quiz_to_disk(filename: str, quiz_data: List[Dict]):
    """Saves a quiz to a JSON file."""
    # Clean filename to be safe
    safe_name = "".join([c for c in filename if c.isalpha() or c.isdigit() or c in " ._-"])
    file_path = os.path.join(QUIZ_STORAGE_DIR, f"{safe_name}.json")
    
    with open(file_path, "w", encoding="utf-8") as f:
        json.dump(quiz_data, f, indent=2)
    return safe_name

def list_saved_quizzes():
    """Returns a list of saved quiz filenames."""
    files = glob.glob(os.path.join(QUIZ_STORAGE_DIR, "*.json"))
    # Return just the filenames without extension
    return [os.path.splitext(os.path.basename(f))[0] for f in files]

def load_quiz_from_disk(filename: str):
    """Loads a quiz by filename."""
    safe_name = "".join([c for c in filename if c.isalpha() or c.isdigit() or c in " ._-"])
    file_path = os.path.join(QUIZ_STORAGE_DIR, f"{safe_name}.json")
    
    if not os.path.exists(file_path):
        return None
        
    with open(file_path, "r", encoding="utf-8") as f:
        return json.load(f)